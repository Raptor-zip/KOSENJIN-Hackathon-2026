#!/usr/bin/env python3
"""Generate NEMUKE BUSTER VC-pitch PPTX (11 slides) — with user validation + Discord growth."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# === Colors ===
DARK_INK = RGBColor(0x1F, 0x29, 0x37)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
EMPHASIS_RED = RGBColor(0xDC, 0x26, 0x26)
SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81)
TEAL = RGBColor(0x08, 0x91, 0xB2)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DISCORD_PURPLE = RGBColor(0x5B, 0x65, 0xEA)

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
CARD_BG = RGBColor(0xF9, 0xFA, 0xFB)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 11
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "presentation.pptx")


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def tb(slide, l, t, w, h, text, sz=18, color=DARK_INK, bold=False,
       align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def tb_rich(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    """Add textbox with individually styled lines: [(text, sz, color, bold), ...]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, sz, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(8)
    return box


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def circle(slide, l, t, sz, fill, text="", txt_color=WHITE, txt_sz=16):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(txt_sz)
        tf.paragraphs[0].font.color.rgb = txt_color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
    return s


def accent_bar(slide):
    rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)


def section_label(slide, text):
    tb(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
       text, sz=12, color=ACCENT_BLUE, bold=True)


def heading(slide, text, y=Inches(0.8)):
    tb(slide, Inches(0.8), y, Inches(11.5), Inches(0.8),
       text, sz=38, color=DARK_INK, bold=True)
    rect(slide, Inches(0.8), y + Inches(0.75), Inches(2), Pt(4), ACCENT_BLUE)


def page_num(slide, num):
    tb(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.3),
       f"{num}/{TOTAL_SLIDES}", sz=11, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT)


def quote_card(slide, x, y, w, h, quote, name, color, bg):
    """Testimonial card with quote and attribution."""
    rrect(slide, x, y, w, h, bg, color)
    # Left accent bar
    rect(slide, x, y, Inches(0.06), h, color)
    # Quote mark
    tb(slide, x + Inches(0.2), y + Inches(0.05), Inches(0.5), Inches(0.5),
       '"', sz=40, color=color, bold=True)
    # Quote text
    tb(slide, x + Inches(0.5), y + Inches(0.2), w - Inches(0.8), h - Inches(0.7),
       quote, sz=17, color=DARK_INK)
    # Name
    tb(slide, x + Inches(0.5), y + h - Inches(0.45), w - Inches(0.8), Inches(0.35),
       f"— {name}", sz=14, color=color, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ================================================================
    # SLIDE 1: TITLE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BG)
    accent_bar(s)

    tb(s, Inches(1.2), Inches(1.5), Inches(11), Inches(1.2),
       "NEMUKE BUSTER", sz=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(1.2), Inches(2.8), Inches(11), Inches(0.6),
       "\u301c\u7720\u6c17\u30d0\u30b9\u30bf\u30fc\u301c", sz=32, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    rect(s, Inches(3.5), Inches(3.8), Inches(6.3), Pt(1), ACCENT_BLUE)
    tb(s, Inches(1.2), Inches(4.0), Inches(11), Inches(0.6),
       "\u7720\u6c17\u3092\u691c\u77e5\u3057\u3001\u904b\u52d5\u3067\u899a\u9192\u3055\u305b\u308b\u3002\u30d6\u30e9\u30a6\u30b6\u3060\u3051\u3067\u3002",
       sz=22, color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER)

    tb(s, Inches(1.2), Inches(5.2), Inches(11), Inches(0.5),
       "Soma Kaibuchi  /  Yuya Tokumitsu", sz=18,
       color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    rrect(s, Inches(4.2), Inches(6.0), Inches(5), Inches(0.45), DARK_BG2, ACCENT_BLUE)
    tb(s, Inches(4.2), Inches(6.0), Inches(5), Inches(0.45),
       "KOSENJIN Hackathon 2026", sz=13, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    page_num(s, 1)

    # ================================================================
    # SLIDE 2: PROBLEM
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "PROBLEM")
    heading(s, '\u5c45\u7720\u308a\u306f\u300c\u3042\u308b\u3042\u308b\u300d\u3067\u306f\u306a\u3044\u3002\u793e\u4f1a\u554f\u984c\u3060\u3002')

    stats = [
        ("OECD\u6700\u4e0b\u4f4d", "\u65e5\u672c\u306e\u7761\u7720\u6642\u9593\n\u5148\u9032\u56fd33\u30ab\u56fd\u4e2d", EMPHASIS_RED, RED_BG),
        ("70%", "20\u4ee3\u306e7\u5272\u304c\n7\u6642\u9593\u672a\u6e80\u306e\u7761\u7720", ACCENT_BLUE, BLUE_BG),
        ("$136B", "\u5c45\u7720\u308a\u306b\u3088\u308b\n\u751f\u7523\u6027\u640d\u5931\uff08\u7c73\u56fd/\u5e74\uff09", PURPLE, PURPLE_BG),
        ("6,400\u4eba", "\u5c45\u7720\u308a\u904b\u8ee2\u306b\u3088\u308b\n\u6b7b\u4ea1\u8005\u6570\uff08\u7c73\u56fd/\u5e74\uff09", EMPHASIS_RED, RED_BG),
    ]

    cw, ch = Inches(2.8), Inches(2.2)
    for i, (big, desc, color, bg) in enumerate(stats):
        x = Inches(0.6) + i * (cw + Inches(0.25))
        rrect(s, x, Inches(2.0), cw, ch, bg, color)
        tb(s, x, Inches(2.2), cw, Inches(0.8),
           big, sz=36, color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.15), Inches(3.1), cw - Inches(0.3), Inches(1.0),
           desc, sz=15, color=DARK_INK, align=PP_ALIGN.CENTER)

    rrect(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.2), RED_BG, EMPHASIS_RED)
    tb_rich(s, Inches(1.2), Inches(4.9), Inches(10.8), Inches(1.0), [
        ("\u65e2\u5b58\u306e\u89e3\u6c7a\u7b56\u306f\u300c\u30a2\u30e9\u30fc\u30e0\u3092\u9e23\u3089\u3059\u3060\u3051\u300d\u3002", 22, EMPHASIS_RED, True),
        ("\u6b62\u3081\u3066\u3001\u307e\u305f\u5bdd\u308b\u3002\u3053\u308c\u304c\u73fe\u5b9f\u3002", 22, EMPHASIS_RED, True),
    ], align=PP_ALIGN.CENTER)

    tb(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.3),
       "Sources: OECD 2021, NSC, Sleep Foundation, Japan Ministry of Health",
       sz=10, color=MEDIUM_GRAY)
    page_num(s, 2)

    # ================================================================
    # SLIDE 3: SOLUTION
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "SOLUTION")
    heading(s, 'NEMUKE BUSTER\u306f\u300c\u691c\u77e5\u300d\u3067\u7d42\u308f\u3089\u306a\u3044\u3002\u300c\u899a\u9192\u300d\u307e\u3067\u5b8c\u7d50\u3059\u308b\u3002')

    flow = [
        ("1", "\u691c\u77e5", "\u30ab\u30e1\u30e9\u3067EAR\n(Eye Aspect Ratio)\n\u3092\u30ea\u30a2\u30eb\u30bf\u30a4\u30e0\u8a08\u6e2c", ACCENT_BLUE),
        ("2", "\u8b66\u544a", "3\u79d2\u9593\u306e\u9589\u773c\u3067\n\u30a2\u30e9\u30fc\u30e0\u767a\u52d5", AMBER),
        ("3", "\u899a\u9192", "\u30b9\u30af\u30ef\u30c3\u30c85\u56de\u306e\n\u5b8c\u904a\u3067\u89e3\u9664\n\uff08\u30dd\u30fc\u30baAI\u304c\u5224\u5b9a\uff09", SUCCESS_GREEN),
    ]

    fw, fh = Inches(3.2), Inches(2.4)
    for i, (num, title, desc, color) in enumerate(flow):
        x = Inches(0.8) + i * (fw + Inches(0.6))
        rrect(s, x, Inches(2.0), fw, fh, WHITE, color)
        rect(s, x, Inches(2.0), fw, Inches(0.06), color)
        circle(s, x + fw / 2 - Inches(0.3), Inches(2.2), Inches(0.55), color, num)
        tb(s, x, Inches(2.85), fw, Inches(0.4),
           title, sz=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(3.3), fw - Inches(0.4), Inches(1.0),
           desc, sz=15, color=DARK_INK, align=PP_ALIGN.CENTER)
        if i < 2:
            ax = x + fw
            tb(s, ax, Inches(2.8), Inches(0.6), Inches(0.5),
               "\u2192", sz=32, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    diffs = [
        ("\u30a4\u30f3\u30b9\u30c8\u30fc\u30eb\u4e0d\u8981", "\u30d6\u30e9\u30a6\u30b6\u3060\u3051\u3067\u52d5\u304f", ACCENT_BLUE, BLUE_BG),
        ("\u30b5\u30fc\u30d0\u30fc\u4e0d\u8981", "\u5168\u51e6\u7406\u304c\u30c7\u30d0\u30a4\u30b9\u4e0a", TEAL, GREEN_BG),
        ("\u30d7\u30e9\u30a4\u30d0\u30b7\u30fc\u5b8c\u5168\u4fdd\u8b77", "\u30c7\u30fc\u30bf\u304c\u30c7\u30d0\u30a4\u30b9\u304b\u3089\u51fa\u306a\u3044", PURPLE, PURPLE_BG),
    ]

    for i, (title, desc, color, bg) in enumerate(diffs):
        x = Inches(0.8) + i * (Inches(3.6) + Inches(0.4))
        rrect(s, x, Inches(4.9), Inches(3.6), Inches(0.9), bg, color)
        tb(s, x + Inches(0.2), Inches(4.95), Inches(3.2), Inches(0.35),
           title, sz=16, color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(5.35), Inches(3.2), Inches(0.35),
           desc, sz=13, color=DARK_INK, align=PP_ALIGN.CENTER)

    page_num(s, 3)

    # ================================================================
    # SLIDE 4: LIVE DEMO
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BG)
    accent_bar(s)
    section_label(s, "LIVE DEMO")

    tb(s, Inches(1.2), Inches(1.2), Inches(11), Inches(1.0),
       "\u30c7\u30e2 \u2014 \u5b9f\u969b\u306b\u4f53\u9a13\u3057\u3066\u304f\u3060\u3055\u3044",
       sz=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    demo = [
        ("START", "\u30d6\u30e9\u30a6\u30b6\u3067URL\u3092\u958b\u304f\u3060\u3051", ACCENT_BLUE),
        ("MONITOR", "\u9854\u30e9\u30f3\u30c9\u30de\u30fc\u30af\u8868\u793a + \u899a\u9192\u30ec\u30d9\u30eb\u30d0\u30fc", TEAL),
        ("DETECT", "\u76ee\u3092\u9589\u3058\u30663\u79d2 \u2192 \u30a2\u30e9\u30fc\u30e0\u767a\u52d5 + Discord\u306b\u6652\u3089\u3057\u6295\u7a3f", EMPHASIS_RED),
        ("SQUAT", "\u30b9\u30af\u30ef\u30c3\u30c85\u56de\u3067\u30ab\u30a6\u30f3\u30c8 \u2192 \u89e3\u9664", SUCCESS_GREEN),
        ("RETURN", "\u899a\u9192\u72b6\u614b\u3067\u76e3\u8996\u30e2\u30fc\u30c9\u306b\u5fa9\u5e30", ACCENT_BLUE),
    ]

    for i, (label, desc, color) in enumerate(demo):
        y = Inches(2.8) + i * Inches(0.85)
        rrect(s, Inches(2.0), y, Inches(1.4), Inches(0.65), color)
        tb(s, Inches(2.0), y + Inches(0.08), Inches(1.4), Inches(0.5),
           label, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, Inches(3.8), y + Inches(0.08), Inches(8), Inches(0.5),
           desc, sz=20, color=RGBColor(0xE5, 0xE7, 0xEB))

    tb(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.4),
       "\u5be9\u67fb\u54e1\u306e\u7686\u3055\u3093\u3082\u3001\u4eca\u591c\u8a66\u305b\u307e\u3059\u3002URL\u3092\u958b\u304f\u3060\u3051\u3067\u3059\u3002",
       sz=18, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    page_num(s, 4)

    # ================================================================
    # SLIDE 5: USER VALIDATION (NEW!)
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "VALIDATION")
    heading(s, "\u5b9f\u969b\u306b\u4f7f\u3063\u305f\u4eba\u306e\u58f0")

    # User 1: Arakawa
    quote_card(s,
               Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5),
               "\u30a2\u30a4\u30c7\u30a2\u304c\u304a\u3082\u3057\u308d\u3044\u3002\n\u30b9\u30af\u30ef\u30c3\u30c8\u3001\u3044\u306d\u3080\u308a\u306e\u691c\u77e5\u3082\u30b9\u30e0\u30fc\u30ba\u3002\nDiscord\u306b\u52dd\u624b\u306b\u6652\u3089\u3055\u308c\u308b\u306e\u304c\u9762\u767d\u3044\u3002",
               "\u8352\u5ddd\u3055\u3093\uff08\u9ad8\u5c02\u751f\uff09",
               ACCENT_BLUE, BLUE_BG)

    # User 2: Kajiyama
    quote_card(s,
               Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.5),
               "\u3051\u3060\u308b\u3055\u304c\u306a\u304f\u306a\u3063\u305f\u3002\u3057\u3083\u304d\u3063\u3068\u3057\u305f\u3002\n\u30b9\u30af\u30ef\u30c3\u30c8\u306f\u5168\u8eab\u904b\u52d5\u3067\u3001\n\u7720\u6c17\u899a\u307e\u3057\u306b\u3044\u3044\u3002",
               "\u68b6\u5c71\u3055\u3093\uff08\u9ad8\u5c02\u751f\uff09",
               SUCCESS_GREEN, GREEN_BG)

    # Key findings
    tb(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4),
       "\u30e6\u30fc\u30b6\u30fc\u30c6\u30b9\u30c8\u7d50\u679c", sz=18, color=DARK_INK, bold=True)

    findings = [
        ("\u2714 2\u4eba\u4e2d2\u4eba\u304c\u300c\u899a\u9192\u52b9\u679c\u3042\u308a\u300d\u3068\u56de\u7b54", SUCCESS_GREEN),
        ("\u2714 \u5c45\u7720\u308a\u691c\u77e5\u30fb\u30b9\u30af\u30ef\u30c3\u30c8\u5224\u5b9a\u3068\u3082\u306b\u300c\u30b9\u30e0\u30fc\u30ba\u300d\u3068\u8a55\u4fa1", ACCENT_BLUE),
        ("\u2714 Discord\u6652\u3089\u3057\u6a5f\u80fd\u304c\u300c\u9762\u767d\u3044\u300d\u2192 \u30d0\u30a4\u30e9\u30eb\u306e\u7a2e", DISCORD_PURPLE),
    ]

    for i, (text, color) in enumerate(findings):
        y = Inches(5.4) + i * Inches(0.42)
        tb(s, Inches(1.2), y, Inches(10), Inches(0.35),
           text, sz=17, color=color, bold=True)

    # Feature requests = product-market fit signal
    rrect(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.45), AMBER_BG, AMBER)
    tb(s, Inches(1.2), Inches(6.72), Inches(10.8), Inches(0.4),
       "\u6539\u5584\u8981\u671b = PMF\u306e\u30b5\u30a4\u30f3: \u300c\u56de\u6570\u8abf\u6574\u300d\u300c\u30b9\u30c8\u30ec\u30c3\u30c1\u8ffd\u52a0\u300d\u300c\u79d2\u6570\u5909\u66f4\u300d\u2192 \u3082\u3063\u3068\u4f7f\u3044\u305f\u3044\u3068\u3044\u3046\u58f0",
       sz=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    page_num(s, 5)

    # ================================================================
    # SLIDE 6: WHY NOW
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "WHY NOW")
    heading(s, "\u30a8\u30c3\u30b8AI\u306e\u7206\u767a\u7684\u6210\u9577\u304c\u3001\u3053\u306e\u30d7\u30ed\u30c0\u30af\u30c8\u3092\u53ef\u80fd\u306b\u3057\u305f")

    trends = [
        ("$23B \u2192 $197B", "\u30a8\u30c3\u30b8AI\u5e02\u5834\uff082034\u5e74\uff09\nCAGR 23.8%", ACCENT_BLUE, BLUE_BG),
        ("WebGPU", "2025\u5e74\u3001\u4e3b\u8981\u30d6\u30e9\u30a6\u30b6\u5168\u5bfe\u5fdc\n\u30d6\u30e9\u30a6\u30b6\u3067GPU\u63a8\u8ad6\u304c\u73fe\u5b9f\u306b", TEAL, GREEN_BG),
        ("91%", "\u30e6\u30fc\u30b6\u30fc\u304c\u30aa\u30f3\u30c7\u30d0\u30a4\u30b9AI\n\u51e6\u7406\u3092\u597d\u3080\u3068\u56de\u7b54", PURPLE, PURPLE_BG),
        ("60 FPS", "MediaPipe\u3067\u9854468\u70b9+\u4f5333\u70b9\n\u30ea\u30a2\u30eb\u30bf\u30a4\u30e0\u691c\u77e5", SUCCESS_GREEN, GREEN_BG),
    ]

    for i, (big, desc, color, bg) in enumerate(trends):
        x = Inches(0.6) + i * (Inches(2.8) + Inches(0.25))
        rrect(s, x, Inches(2.0), Inches(2.8), Inches(2.4), bg, color)
        tb(s, x, Inches(2.25), Inches(2.8), Inches(0.7),
           big, sz=32, color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.15), Inches(3.1), Inches(2.5), Inches(1.1),
           desc, sz=14, color=DARK_INK, align=PP_ALIGN.CENTER)

    rrect(s, Inches(2.5), Inches(4.9), Inches(8.3), Inches(0.6), BLUE_BG, ACCENT_BLUE)
    tb(s, Inches(2.5), Inches(4.9), Inches(8.3), Inches(0.6),
       "2\u5e74\u524d\u306f\u4e0d\u53ef\u80fd\u3060\u3063\u305f\u3002\u4eca\u3060\u304b\u3089\u3053\u305d\u3067\u304d\u308b\u3002",
       sz=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.3),
       "Sources: market.us, Protecto AI Privacy Report 2025, Google MediaPipe",
       sz=10, color=MEDIUM_GRAY)
    page_num(s, 6)

    # ================================================================
    # SLIDE 7: MARKET (TAM/SAM/SOM)
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "MARKET")
    heading(s, "$8.9B\u306e\u5e02\u5834\u306b\u3001\u30bd\u30d5\u30c8\u30a6\u30a7\u30a2\u3067\u5207\u308a\u8fbc\u3080")

    market = [
        ("TAM", "$27.2B", "\u5c45\u7720\u308a\u691c\u77e5\u30b7\u30b9\u30c6\u30e0\n\u5e02\u5834\u5168\u4f53\uff082034\u5e74\uff09", ACCENT_BLUE, BLUE_BG),
        ("SAM", "$3.5B", "\u30bd\u30d5\u30c8\u30a6\u30a7\u30a2\u30d9\u30fc\u30b9DMS +\n\u5b66\u7fd2\u30fb\u52b4\u50cd\u751f\u7523\u6027\u30c4\u30fc\u30eb", TEAL, GREEN_BG),
        ("SOM", "$50M", "\u65e5\u672c\u306e\u5b66\u751f\u30fb\n\u30ea\u30e2\u30fc\u30c8\u30ef\u30fc\u30ab\u30fc\u30fb\n\u4e2d\u5c0f\u904b\u9001\u696d", SUCCESS_GREEN, GREEN_BG),
    ]

    mw, mh = Inches(3.5), Inches(2.8)
    for i, (label, amount, desc, color, bg) in enumerate(market):
        x = Inches(0.6) + i * (mw + Inches(0.4))
        rrect(s, x, Inches(2.0), mw, mh, bg, color)
        bw = Inches(1.0)
        rrect(s, x + mw / 2 - bw / 2, Inches(2.15), bw, Inches(0.4), color)
        tb(s, x + mw / 2 - bw / 2, Inches(2.15), bw, Inches(0.4),
           label, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x, Inches(2.7), mw, Inches(0.7),
           amount, sz=40, color=color, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(3.5), mw - Inches(0.4), Inches(1.2),
           desc, sz=15, color=DARK_INK, align=PP_ALIGN.CENTER)

    rrect(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.9), RED_BG, EMPHASIS_RED)
    tb(s, Inches(1.2), Inches(5.4), Inches(10.8), Inches(0.7),
       "\u65e2\u5b58\u30d7\u30ec\u30a4\u30e4\u30fc\u306f\u5168\u3066\u30cf\u30fc\u30c9\u30a6\u30a7\u30a2\u4f9d\u5b58\uff08$500\u301c$5,000/\u53f0\uff09\u3002NEMUKE BUSTER\u306f\u9650\u754c\u8cbb\u7528\u307b\u307c\u30bc\u30ed\u3002",
       sz=20, color=EMPHASIS_RED, bold=True, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3),
       "Sources: market.us, GM Insights, Grand View Research",
       sz=10, color=MEDIUM_GRAY)
    page_num(s, 7)

    # ================================================================
    # SLIDE 8: COMPETITIVE ADVANTAGE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "ADVANTAGE")
    heading(s, "\u306a\u305cNEMUKE BUSTER\u304c\u52dd\u3066\u308b\u306e\u304b")

    headers = ["", "\u8eca\u8f09DMS", "\u5c02\u7528\u30c7\u30d0\u30a4\u30b9", "\u30a2\u30e9\u30fc\u30e0\u30a2\u30d7\u30ea", "NEMUKE BUSTER"]
    rows = [
        headers,
        ["\u30b3\u30b9\u30c8", "$500\u301c5,000", "$200\u301c500", "\u7121\u6599", "\u7121\u6599 / Freemium"],
        ["\u30a4\u30f3\u30b9\u30c8\u30fc\u30eb", "\u8eca\u4e21\u7d44\u8fbc", "\u30c7\u30d0\u30a4\u30b9\u8cfc\u5165", "\u30a2\u30d7\u30eaDL", "URL\u3092\u958b\u304f\u3060\u3051"],
        ["\u899a\u9192\u30a2\u30af\u30b7\u30e7\u30f3", "\u306a\u3057", "\u306a\u3057", "\u306a\u3057", "\u30b9\u30af\u30ef\u30c3\u30c8\u5224\u5b9a"],
        ["\u30d7\u30e9\u30a4\u30d0\u30b7\u30fc", "\u30af\u30e9\u30a6\u30c9\u9001\u4fe1", "\u30af\u30e9\u30a6\u30c9\u9001\u4fe1", "\u30ed\u30fc\u30ab\u30eb", "\u5b8c\u5168\u30ed\u30fc\u30ab\u30eb"],
        ["\u30bd\u30fc\u30b7\u30e3\u30eb\u6a5f\u80fd", "\u306a\u3057", "\u306a\u3057", "\u306a\u3057", "Discord\u6652\u3089\u3057"],
        ["\u5229\u7528\u5834\u6240", "\u8eca\u5185\u306e\u307f", "\u56fa\u5b9a\u5834\u6240", "\u3069\u3053\u3067\u3082", "\u3069\u3053\u3067\u3082"],
    ]

    n_rows = len(rows)
    n_cols = len(headers)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(2.0),
                                    Inches(11.5), Inches(n_rows * 0.55))
    tbl = tbl_shape.table

    col_widths = [Inches(2.0), Inches(2.1), Inches(2.1), Inches(2.1), Inches(3.2)]
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = txt
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(13)
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                elif ci == 4:
                    p.font.color.rgb = ACCENT_BLUE
                    p.font.bold = True
                else:
                    p.font.color.rgb = DARK_INK

            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_BLUE if ci == 4 else DARK_INK
            elif ci == 4:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_BG if ri % 2 == 0 else RGBColor(0xDB, 0xEA, 0xFE)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG

    page_num(s, 8)

    # ================================================================
    # SLIDE 9: TECH MOAT
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "TECHNOLOGY")
    heading(s, "\u30d6\u30e9\u30a6\u30b6\u30cd\u30a4\u30c6\u30a3\u30d6AI\u3068\u3044\u3046\u6280\u8853\u7684\u512a\u4f4d")

    tech = [
        ("React 19 + TypeScript + Vite", "\u30e2\u30c0\u30f3Web\u6280\u8853\u30b9\u30bf\u30c3\u30af", ACCENT_BLUE),
        ("MediaPipe FaceLandmarker", "468\u70b9\u306e\u9854\u7279\u5fb4\u70b9\u3067EAR\u7b97\u51fa", TEAL),
        ("MediaPipe PoseLandmarker", "33\u70b9\u306e\u4f53\u95a2\u7bc0\u3067\u30b9\u30af\u30ef\u30c3\u30c8\u5224\u5b9a", SUCCESS_GREEN),
        ("\u8aa4\u691c\u77e5\u9632\u6b62\u30a2\u30eb\u30b4\u30ea\u30ba\u30e0", "3\u30d5\u30ec\u30fc\u30e0\u9023\u7d9a\u4e00\u81f4 + 800ms\u30af\u30fc\u30eb\u30c0\u30a6\u30f3", EMPHASIS_RED),
        ("Discord Webhook\u9023\u643a", "\u5c45\u7720\u308a\u30b9\u30af\u30b7\u30e7\u3092\u81ea\u52d5\u6295\u7a3f \u2192 \u30bd\u30fc\u30b7\u30e3\u30eb\u30d7\u30ec\u30c3\u30b7\u30e3\u30fc\u3067\u899a\u9192\u4fc3\u9032", DISCORD_PURPLE),
    ]

    for i, (title, desc, color) in enumerate(tech):
        y = Inches(2.0) + i * Inches(0.85)
        rrect(s, Inches(1.0), y, Inches(10.5), Inches(0.7), CARD_BG, LIGHT_GRAY)
        rect(s, Inches(1.0), y, Inches(0.08), Inches(0.7), color)
        circle(s, Inches(1.3), y + Inches(0.08), Inches(0.5), color, str(i + 1), WHITE, 14)
        tb(s, Inches(2.1), y + Inches(0.05), Inches(4.5), Inches(0.3),
           title, sz=17, color=DARK_INK, bold=True)
        tb(s, Inches(2.1), y + Inches(0.35), Inches(9), Inches(0.3),
           desc, sz=14, color=MEDIUM_GRAY)

    rrect(s, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6), BLUE_BG, ACCENT_BLUE)
    tb(s, Inches(1.5), Inches(6.35), Inches(10), Inches(0.5),
       "\u30b5\u30fc\u30d0\u30fc\u30b3\u30b9\u30c8 $0\u3002100\u4e07\u30e6\u30fc\u30b6\u30fc\u3067\u3082\u904b\u7528\u30b3\u30b9\u30c8\u304c\u5897\u3048\u306a\u3044\u3002",
       sz=20, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    page_num(s, 9)

    # ================================================================
    # SLIDE 10: BUSINESS MODEL & GROWTH
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    section_label(s, "GROWTH")
    heading(s, "Discord\u6652\u3089\u3057\u3067\u30d0\u30a4\u30e9\u30eb\u3001B2B\u3067\u53ce\u76ca\u5316")

    phases = [
        ("Phase 1", "Now", "\u7121\u6599\u3067\u30d0\u30a4\u30e9\u30eb\u7372\u5f97",
         ["\u57fa\u672c\u6a5f\u80fd\u3092\u7121\u6599\u63d0\u4f9b",
          "Discord\u6652\u3089\u3057 \u2192 \u53cb\u4eba\u304c\u8a66\u3059",
          "\u300cNEMUKE BUSTER\u3055\u308c\u305f\u300d\u304cSNS\u3067\u62e1\u6563"],
         ACCENT_BLUE, BLUE_BG),
        ("Phase 2", "6\u30f6\u6708\u5f8c", "\u30d7\u30ec\u30df\u30a2\u30e0\u5316",
         ["\u899a\u9192\u30c7\u30fc\u30bf\u5206\u6790\u30c0\u30c3\u30b7\u30e5\u30dc\u30fc\u30c9",
          "\u30ab\u30b9\u30bf\u30e0\u904b\u52d5\u30e1\u30cb\u30e5\u30fc",
          "PWA\u5316\uff08\u30aa\u30d5\u30e9\u30a4\u30f3\u5bfe\u5fdc\uff09"],
         TEAL, GREEN_BG),
        ("Phase 3", "1\u5e74\u5f8c", "B2B\u5c55\u958b",
         ["\u904b\u9001\u4f1a\u793e\u5411\u3051\u30c9\u30e9\u30a4\u30d0\u30fc\u7ba1\u7406",
          "\u4f01\u696d\u5411\u3051\u751f\u7523\u6027\u30c4\u30fc\u30eb",
          "\u6559\u80b2\u6a5f\u95a2\u5411\u3051\u30e9\u30a4\u30bb\u30f3\u30b9"],
         PURPLE, PURPLE_BG),
    ]

    pw, ph_val = Inches(3.5), Inches(3.4)
    for i, (phase, when, title, bullets, color, bg) in enumerate(phases):
        x = Inches(0.7) + i * (pw + Inches(0.45))
        rrect(s, x, Inches(2.0), pw, ph_val, bg, color)
        rect(s, x, Inches(2.0), pw, Inches(0.06), color)

        bw = Inches(1.8)
        rrect(s, x + pw / 2 - bw / 2, Inches(2.15), bw, Inches(0.4), color)
        tb(s, x + pw / 2 - bw / 2, Inches(2.15), bw, Inches(0.4),
           f"{phase}  |  {when}", sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        tb(s, x + Inches(0.2), Inches(2.7), pw - Inches(0.4), Inches(0.4),
           title, sz=18, color=color, bold=True, align=PP_ALIGN.CENTER)

        for j, bullet in enumerate(bullets):
            by = Inches(3.3) + j * Inches(0.55)
            tb(s, x + Inches(0.25), by, pw - Inches(0.5), Inches(0.5),
               f"\u2022 {bullet}", sz=13, color=DARK_INK)

        if i < 2:
            ax = x + pw
            tb(s, ax, Inches(3.2), Inches(0.45), Inches(0.5),
               "\u2192", sz=32, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Growth flywheel
    rrect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8), RGBColor(0xED, 0xEF, 0xFE), DISCORD_PURPLE)
    tb(s, Inches(1.2), Inches(5.85), Inches(10.8), Inches(0.7),
       "\u25b6 \u5c45\u7720\u308a\u691c\u77e5 \u2192 Discord\u306b\u30b9\u30af\u30b7\u30e7\u6295\u7a3f \u2192 \u53cb\u4eba\u304c\u7b11\u3046 \u2192 \u53cb\u4eba\u3082\u4f7f\u3046 \u2192 \u30ea\u30d4\u30fc\u30c8",
       sz=18, color=DISCORD_PURPLE, bold=True, align=PP_ALIGN.CENTER)

    page_num(s, 10)

    # ================================================================
    # SLIDE 11: CLOSING
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, DARK_BG)
    accent_bar(s)

    tb(s, Inches(1.2), Inches(1.0), Inches(11), Inches(1.0),
       "NEMUKE BUSTER", sz=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(5.0), Inches(2.2), Inches(3.3), Pt(2), ACCENT_BLUE)

    tb(s, Inches(1.2), Inches(2.5), Inches(11), Inches(0.6),
       "$27B\u5e02\u5834\u306b\u3001\u30bc\u30ed\u30b3\u30b9\u30c8\u306e\u30bd\u30d5\u30c8\u30a6\u30a7\u30a2\u3067\u6311\u3080\u3002",
       sz=26, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    tb(s, Inches(1.2), Inches(3.3), Inches(11), Inches(0.5),
       "\u7720\u6c17\u3092\u691c\u77e5\u3057\u3001\u904b\u52d5\u3067\u899a\u9192\u3055\u305b\u308b\u3002\u30cf\u30fc\u30c9\u30a6\u30a7\u30a2\u4e0d\u8981\u3002\u30d6\u30e9\u30a6\u30b6\u3060\u3051\u3067\u3002",
       sz=20, color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER)

    # QR code placeholder
    qr_size = Inches(2.0)
    qr_x = SLIDE_WIDTH / 2 - qr_size / 2
    rrect(s, qr_x, Inches(4.0), qr_size, qr_size, DARK_BG2, ACCENT_BLUE)
    tb(s, qr_x, Inches(4.3), qr_size, Inches(0.5),
       "[ QR CODE ]", sz=18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    tb(s, qr_x, Inches(4.8), qr_size, Inches(0.5),
       "\u2191 \u3053\u3053\u306bQR\u3092\u8cbc\u308b", sz=12, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    tb(s, Inches(1.2), Inches(6.2), Inches(11), Inches(0.5),
       "\u4eca\u3059\u3050\u8a66\u305b\u307e\u3059\u3002", sz=24, color=ACCENT_BLUE,
       bold=True, align=PP_ALIGN.CENTER)

    tb(s, Inches(1.2), Inches(6.7), Inches(11), Inches(0.5),
       "\u3042\u308a\u304c\u3068\u3046\u3054\u3056\u3044\u307e\u3057\u305f\uff01", sz=28, color=WHITE,
       bold=True, align=PP_ALIGN.CENTER)

    page_num(s, 11)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
